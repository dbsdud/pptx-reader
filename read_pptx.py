#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint Reader Skill for Claude
Extracts text content from .pptx files
"""

import sys
import json
from pathlib import Path
import io

# Set UTF-8 encoding for output
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

try:
    from pptx import Presentation
except ImportError:
    print(json.dumps({
        "error": "python-pptx library not installed",
        "message": "Please install it with: pip install python-pptx"
    }))
    sys.exit(1)


def extract_text_from_slide(slide):
    """Extract all text from a slide"""
    text_content = []

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text_content.append(shape.text.strip())

    return "\n".join(text_content)


def read_pptx(file_path):
    """Read PowerPoint file and extract content"""
    try:
        prs = Presentation(file_path)

        result = {
            "file_path": str(file_path),
            "total_slides": len(prs.slides),
            "slides": []
        }

        for idx, slide in enumerate(prs.slides, 1):
            slide_text = extract_text_from_slide(slide)

            slide_info = {
                "slide_number": idx,
                "content": slide_text
            }

            result["slides"].append(slide_info)

        return result

    except FileNotFoundError:
        return {"error": f"File not found: {file_path}"}
    except Exception as e:
        return {"error": f"Error reading PowerPoint: {str(e)}"}


def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "No file path provided"}))
        sys.exit(1)

    file_path = sys.argv[1]
    result = read_pptx(file_path)

    # Format output for Claude
    if "error" in result:
        print(json.dumps(result, indent=2, ensure_ascii=False))
    else:
        print(f"\n# PowerPoint Content: {Path(file_path).name}\n")
        print(f"Total Slides: {result['total_slides']}\n")
        print("=" * 80 + "\n")

        for slide in result["slides"]:
            print(f"## Slide {slide['slide_number']}")
            print("-" * 80)
            print(slide['content'])
            print("\n" + "=" * 80 + "\n")


if __name__ == "__main__":
    main()
